"""Extracteurs réels de `core.enrich` sur des fichiers générés à la volée.

Les libs bureautiques (pypdf, python-docx, openpyxl, python-pptx) sont déjà des
dépendances du moteur : on génère de vrais petits fichiers en tmp_path et on
vérifie que chaque extracteur en tire propriétés et extrait de texte. Aucun
binaire n'est committé, aucun réseau.
"""
import pytest

from core.enrich import (
    _extract_docx,
    _extract_pdf,
    _extract_pptx,
    _extract_xlsx,
    enrich_descriptions,
)


@pytest.fixture
def docx_file(tmp_path):
    from docx import Document

    doc = Document()
    doc.core_properties.title = "Budget cantine 2022"
    doc.core_properties.keywords = "cantine, budget"
    doc.add_paragraph("Compte rendu de la commission cantine.")
    path = tmp_path / "rapport.docx"
    doc.save(str(path))
    return path


@pytest.fixture
def xlsx_file(tmp_path):
    from openpyxl import Workbook

    wb = Workbook()
    wb.properties.title = "Effectifs rentrée"
    wb.active.title = "Effectifs 2023"
    wb.create_sheet("Dérogations")
    path = tmp_path / "effectifs.xlsx"
    wb.save(str(path))
    return path


@pytest.fixture
def pptx_file(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.core_properties.subject = "Rentrée scolaire"
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Presentation de la rentree aux parents."
    path = tmp_path / "reunion.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def pdf_file(tmp_path):
    from pypdf import PdfWriter

    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    writer.add_metadata({"/Subject": "Facture traiteur", "/Author": "Mairie"})
    path = tmp_path / "facture.pdf"
    with open(path, "wb") as f:
        writer.write(f)
    return path


def test_extract_docx_props_and_snippet(docx_file):
    fields = _extract_docx(docx_file)
    assert fields["Titre interne"] == "Budget cantine 2022"
    assert fields["Mots-clés"] == "cantine, budget"
    assert "commission cantine" in fields["Extrait"]
    assert "Auteur" not in fields  # l'auteur n'est plus lu (donnée personnelle inutile)


def test_extract_xlsx_props_and_sheets(xlsx_file):
    fields = _extract_xlsx(xlsx_file)
    assert fields["Titre interne"] == "Effectifs rentrée"
    assert fields["Feuilles"] == "Effectifs 2023, Dérogations"
    assert "Auteur" not in fields


def test_extract_pptx_props_and_text(pptx_file):
    fields = _extract_pptx(pptx_file)
    assert fields["Sujet"] == "Rentrée scolaire"
    assert "rentree aux parents" in fields["Extrait"]
    assert "Auteur" not in fields


def test_extract_pdf_metadata(pdf_file):
    fields = _extract_pdf(pdf_file)
    assert fields["Sujet"] == "Facture traiteur"
    assert "Auteur" not in fields  # l'auteur n'est plus lu (donnée personnelle inutile)
    assert "Extrait" not in fields  # page blanche : pas de couche texte


def test_enrich_end_to_end_with_real_files(tmp_path, docx_file, pdf_file):
    import pandas as pd

    df = pd.DataFrame([
        {"File": docx_file.name, "Content.DescriptionLevel": "Item",
         "Content.Title": docx_file.name},
        {"File": pdf_file.name, "Content.DescriptionLevel": "Item",
         "Content.Title": pdf_file.name},
    ])
    out, report = enrich_descriptions(df, tmp_path)
    assert report.enriched == 2
    assert "Budget cantine 2022" in out.loc[0, "Content.Description"]
    assert "Facture traiteur" in out.loc[1, "Content.Description"]
