"""Enrichissement mécanique des descriptions à partir des binaires.

Étape de préparation **optionnelle**, exécutée en local sur la machine de
l'archiviste, **avant** l'audit. Elle lit les fichiers bureautiques référencés
par la colonne `File` du CSV et en extrait des métadonnées **déterministes**
(propriétés du document, mots-clés, premières lignes de texte) qu'elle écrit
dans `Content.Description`. Cette colonne est ensuite consommée telle quelle par
AUD-001 et CLA-001 (option « Inclure la description »).

Particularité assumée : contrairement au reste du moteur — qui ne lit que des
métadonnées de chemin/nom/date —, cette étape **ouvre le contenu** des fichiers.
C'est un choix explicite de l'utilisateur : l'étape est facultative, aucune
donnée ne quitte la machine, et l'appelant (CLI/backend) doit afficher un
message clair sur l'accès aux données avant de la lancer.

Aucun appel LLM, aucune OCR : seuls le texte déjà présent dans le fichier
(couche texte des PDF, corps des .docx, etc.) et les propriétés embarquées sont
exploités. Un PDF scanné sans couche texte ne produit donc rien — c'est normal.
"""
from __future__ import annotations

import hashlib
import re
from collections import Counter
from collections.abc import Mapping
from dataclasses import dataclass, field
from pathlib import Path

import pandas as pd

# Extensions traitées (extraction mécanique de texte/propriétés possible).
# Les autres (jpg, png, zip, …) sont ignorées : pas de texte exploitable
# sans OCR, hors périmètre de l'enrichissement mécanique.
SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx"}

# Colonne portant l'empreinte SHA-256 d'un binaire. Produite par
# `fingerprint_files`, consommée par `core.audit_scan` pour détecter les
# **doublons stricts** (fichiers binairement identiques). Hors RESIP standard —
# colonne d'appoint, non requise par l'export.
FINGERPRINT_COLUMN = "Content.Fingerprint"

# Lecture par blocs pour ne jamais charger un gros binaire entièrement en mémoire.
_HASH_CHUNK = 1 << 20  # 1 MiB

# Longueur par défaut de la description produite (budget tokens en aval).
DEFAULT_MAX_CHARS = 300
_TEXT_SNIPPET_CHARS = 220
# Unités (paragraphes docx, lignes xlsx) de part et d'autre de chaque point
# d'échantillonnage début/milieu/fin (cf. `_sample_indices`/`_windowed_indices`).
_SAMPLE_WINDOW = 3


# ── Rapport ──────────────────────────────────────────────────────────────────

@dataclass
class EnrichReport:
    total_items: int = 0
    enriched: int = 0          # description écrite
    already_filled: int = 0    # description existante préservée (sans --overwrite)
    no_text: int = 0           # fichier lu mais rien d'exploitable (ex. PDF scanné)
    unsupported: int = 0       # extension hors périmètre (jpg, …)
    missing: int = 0           # binaire introuvable sous source_root
    errors: list[str] = field(default_factory=list)

    def summary_lines(self) -> list[str]:
        out = [
            f"{self.enriched} description(s) écrite(s) sur {self.total_items} item(s)",
        ]
        if self.already_filled:
            out.append(f"{self.already_filled} description(s) existante(s) préservée(s)")
        if self.no_text:
            out.append(f"{self.no_text} fichier(s) sans texte exploitable (PDF scanné, vide…)")
        if self.unsupported:
            out.append(f"{self.unsupported} fichier(s) ignoré(s) (format non bureautique)")
        if self.missing:
            out.append(f"{self.missing} binaire(s) introuvable(s) sous la racine")
        if self.errors:
            out.append(f"{len(self.errors)} erreur(s) de lecture")
        return out


@dataclass
class FingerprintReport:
    total_items: int = 0
    hashed: int = 0            # empreinte calculée
    already_hashed: int = 0    # empreinte existante préservée (sans --overwrite)
    missing: int = 0           # binaire introuvable sous source_root
    skipped: int = 0           # ligne sans chemin de fichier exploitable
    errors: list[str] = field(default_factory=list)

    def summary_lines(self) -> list[str]:
        out = [
            f"{self.hashed} empreinte(s) SHA-256 calculée(s) sur {self.total_items} item(s)",
        ]
        if self.already_hashed:
            out.append(f"{self.already_hashed} empreinte(s) existante(s) préservée(s)")
        if self.missing:
            out.append(f"{self.missing} binaire(s) introuvable(s) sous la racine")
        if self.errors:
            out.append(f"{len(self.errors)} erreur(s) de lecture")
        return out


# ── Avertissement d'accès au contenu (source unique CLI ⇄ backend) ────────────

def content_access_notice_lines(source_root: str | Path) -> list[str]:
    """Message d'avertissement de l'étape `enrich` — **source unique** partagée
    par la CLI (`cmd_enrich`) et le backend HTTP (`/enrich`).

    Cette étape est la seule du moteur à **ouvrir le contenu** des fichiers (le
    reste ne lit que des métadonnées de chemin/nom/date). Le message doit être
    affiché avant traitement : c'est le consentement explicite de l'archiviste.
    Rien ne quitte la machine ; aucun appel LLM ; pas d'OCR.
    """
    return [
        "Étape de préparation (facultative) — accès au contenu des fichiers.",
        f"Les binaires de « {source_root} » vont être OUVERTS en local pour en "
        "extraire des métadonnées (propriétés, mots-clés, premières lignes).",
        "Aucune donnée ne quitte la machine ; aucun appel LLM ; pas d'OCR.",
    ]


# ── Utilitaires ──────────────────────────────────────────────────────────────

_EMAIL_RE = re.compile(r"\b[A-Za-z0-9._%+-]+@([A-Za-z0-9.-]+\.[A-Za-z]{2,})\b")


def _redact_emails(text: str) -> str:
    """Réduit une adresse mail à son domaine, avec un repère explicite
    (`jean.dupont@mairie.fr` → `[email]@mairie.fr`) : la partie locale
    identifie une personne (donnée personnelle sans utilité pour la
    classification), le domaine identifie l'organisme émetteur (signal utile
    — typologie/activité — à conserver). Le placeholder `[email]` évite qu'un
    `@mairie.fr` isolé ressemble à un artefact de troncature illisible plutôt
    qu'à une adresse volontairement anonymisée (côté lecteur humain ou LLM)."""
    return _EMAIL_RE.sub(r"[email]@\1", text)


def _clean(text: str | None) -> str:
    """Normalise un fragment : espaces compactés, retours ligne supprimés,
    adresses mail réduites au domaine (`_redact_emails`)."""
    if not text:
        return ""
    cleaned = re.sub(r"\s+", " ", str(text)).strip()
    return _redact_emails(cleaned)


def _resolve(source_root: Path, file_value: str) -> Path:
    """Résout un chemin `File` (séparateurs Windows possibles) sous la racine."""
    parts = [p for p in file_value.replace("\\", "/").split("/") if p not in ("", ".")]
    return source_root.joinpath(*parts)


def _is_noise_title(title: str, filename: str) -> bool:
    """Un « titre interne » égal au nom de fichier ou à un libellé d'outil
    n'apporte rien — on l'écarte pour ne pas polluer la description."""
    t = title.lower().strip()
    if not t:
        return True
    if t == filename.lower() or t == Path(filename).stem.lower():
        return True
    # Titres génériques produits par les suites bureautiques / scanners
    if t.startswith("microsoft word -"):
        return True
    return t in (
        "document", "presentation", "classeur",
        "fichier pdf", "fichier", "untitled", "sans titre",
    )


def _looks_like_text(text: str, min_alpha_ratio: float = 0.6) -> bool:
    """Rejette un fragment qui ne ressemble pas à de la langue naturelle.

    Certains PDF (polices à encodage non standard, sans table ToUnicode) font
    extraire à pypdf des codes bruts illisibles (« 3$<*Y 05'=$1… »). Ce bruit
    serait pire que rien dans Content.Description — on le détecte via le ratio
    de caractères alphabétiques sur les caractères non-espace. Utilisé pour
    l'extrait de corps de texte *et* pour le titre repéré dans le contenu (même
    origine possible : une ligne brute extraite d'un PDF à police illisible).
    """
    non_space = [c for c in text if not c.isspace()]
    if not non_space:
        return False
    alpha_ratio = sum(c.isalpha() for c in non_space) / len(non_space)
    return alpha_ratio >= min_alpha_ratio


def _plausible_snippet(raw: str | None) -> str:
    """Nettoie un extrait de corps de texte et le rejette s'il est illisible
    (cf. `_looks_like_text`)."""
    text = _clean(raw)
    if not text or not _looks_like_text(text):
        return ""
    return text[:_TEXT_SNIPPET_CHARS]


# ── Extraction intelligente : échantillonnage, thème, en-tête/pied de page ────
# Fonctions pures (pas d'accès disque/lib tierce) : le bruit est réduit par
# l'algorithme, pas par un appel LLM — dans l'esprit déterministe du module.

# Mots vides français (articles, pronoms, prépositions, auxiliaires, formules
# de politesse de courrier) — écartés du calcul de fréquence pour ne pas noyer
# les mots réellement récurrents et thématiques. Source unique : également
# consommés par l'agrégation de termes du vrac (core.agt_tools.mots_frequents).
FRENCH_STOPWORDS = frozenset({
    "alors", "au", "aucun", "aucune", "aussi", "autre", "autres", "avant",
    "avec", "avoir", "bon", "car", "ce", "cela", "ces", "cet", "cette",
    "ceux", "chaque", "ci", "comme", "comment", "dans", "de", "des", "du",
    "dedans", "dehors", "depuis", "devrait", "doit", "donc", "dont", "elle",
    "elles", "en", "encore", "entre", "est", "et", "eu", "eux", "fait",
    "faites", "fois", "font", "hors", "ici", "il", "ils", "je", "juste",
    "la", "le", "les", "leur", "leurs", "là", "ma", "maintenant", "mais",
    "mes", "moins", "mon", "même", "mêmes", "ni", "notre", "nos", "nous",
    "nouveaux", "ou", "où", "par", "parce", "parole", "pas", "personnes",
    "peut", "peu", "plupart", "pour", "pourquoi", "quand", "que", "quel",
    "quelle", "quelles", "quels", "qui", "sa", "sans", "se", "ses",
    "seulement", "si", "sien", "son", "sont", "sous", "soyez", "sujet",
    "sur", "ta", "tandis", "tellement", "tels", "tes", "toi", "ton", "tous",
    "tout", "toute", "toutes", "trop", "très", "tu", "un", "une", "vos",
    "votre", "vous", "vu", "ça", "étaient", "état", "étions",
    "été", "être", "avez", "avons", "ai", "as", "a", "ont", "avait",
    "avaient", "suis", "es", "sommes", "êtes", "serai", "seras", "sera",
    "serons", "seront", "étais", "était", "fut", "furent",
    "monsieur", "madame", "mademoiselle", "cordialement", "bonjour",
})

# Mots injectés par nos propres placeholders (`_redact_emails`) : jamais un
# signal thématique, à écarter même s'ils reviennent (plusieurs adresses mail
# dans un même document produiraient sinon un faux "mot récurrent").
_PLACEHOLDER_WORDS = frozenset({"email"})

_WORD_RE = re.compile(r"[A-Za-zÀ-ÖØ-öø-ÿ]{3,}")


def _sample_indices(total: int, want: int = 3) -> list[int]:
    """Positions début/milieu/fin dans `[0, total)`, sans parcourir tout le
    document (pages PDF, paragraphes docx, slides pptx, lignes xlsx)."""
    if total <= 0:
        return []
    if total <= want:
        return list(range(total))
    return sorted({0, total // 2, total - 1})


def _windowed_indices(anchors: list[int], total: int, window: int) -> list[int]:
    """Étend chaque point d'ancrage en une petite fenêtre contiguë bornée,
    pour capter un peu de matière autour de chaque échantillon plutôt qu'une
    seule unité isolée (ex. un seul paragraphe)."""
    idx: set[int] = set()
    for anchor in anchors:
        lo = max(0, anchor - window)
        hi = min(total - 1, anchor + window)
        idx.update(range(lo, hi + 1))
    return sorted(idx)


def _tokenize(text: str) -> list[str]:
    """Mots en minuscules (lettres accentuées incluses), longueur >= 3."""
    return [w.lower() for w in _WORD_RE.findall(text)]


def _recurring_keywords(text: str, top_n: int = 6) -> str:
    """Mots qui reviennent au moins deux fois dans `text`, hors mots vides.

    Simple fréquence, pas de TF-IDF : un mot mentionné une seule fois n'est
    pas « récurrent » au sens où l'utilisateur l'entend — le seuil >= 2 est
    le filtre anti-bruit. Emails réduits au domaine avant tokenisation : cette
    fonction reçoit du texte brut (pas toujours passé par `_clean` en amont).
    """
    tokens = [
        w for w in _tokenize(_redact_emails(text))
        if w not in FRENCH_STOPWORDS and w not in _PLACEHOLDER_WORDS
    ]
    if not tokens:
        return ""
    counts = Counter(tokens)
    recurring = [word for word, n in counts.most_common() if n >= 2]
    return ", ".join(recurring[:top_n])


def _detect_boilerplate_lines(texts: list[str]) -> set[str]:
    """Lignes identiques sur >= 2 extraits = en-tête/pied de page récurrent.

    Utilisé pour le PDF, qui n'expose aucune API d'en-tête/pied de page
    (contrairement à docx/xlsx/pptx où on lit l'info directement) : on
    l'infère par répétition entre pages échantillonnées.
    """
    if len(texts) < 2:
        return set()
    counts: Counter[str] = Counter()
    for text in texts:
        lines = {_clean(line) for line in text.splitlines() if _clean(line)}
        for line in lines:
            if len(line) <= 150:
                counts[line] += 1
    return {line for line, n in counts.items() if n >= 2}


_TITLE_STYLES = ("title", "heading")


def _pick_content_title(candidates: list[tuple[str, str]], filename: str) -> str:
    """Premier candidat `(texte, style_ou_origine)` non-bruit, en donnant la
    priorité aux styles structurels (« Title »/« Heading* ») sur la simple
    position dans le document."""
    styled = [c for c in candidates if c[1].lower().startswith(_TITLE_STYLES)]
    for text, _origin in (*styled, *candidates):
        cleaned = _clean(text)
        if cleaned and not _is_noise_title(cleaned, filename) and _looks_like_text(cleaned):
            return cleaned
    return ""


# ── Extracteurs par format ───────────────────────────────────────────────────
# Chacun retourne un dict ordonné de champs {label: valeur}. Import paresseux :
# une dépendance manquante n'empêche que son format, pas tout l'enrichissement.

def _extract_pdf(path: Path) -> dict[str, str]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    fields: dict[str, str] = {}
    # DocumentInformation (Mapping) ou dict vide : seul .get() est utilisé.
    meta: Mapping = reader.metadata or {}

    title = _clean(meta.get("/Title"))
    if title and not _is_noise_title(title, path.name):
        fields["Titre interne"] = title

    # Pages début/milieu/fin plutôt que les 2 premières uniquement — le thème
    # peut se trouver n'importe où dans un document long.
    page_indices = _sample_indices(len(reader.pages))
    page_texts: list[str] = []
    for i in page_indices:
        try:
            page_texts.append(reader.pages[i].extract_text() or "")
        except Exception:
            page_texts.append("")

    # En-tête/pied de page : pypdf n'expose rien nativement, on l'infère par
    # récurrence de lignes identiques entre pages échantillonnées.
    boilerplate = _detect_boilerplate_lines(page_texts)
    body_texts = [
        "\n".join(line for line in text.splitlines() if _clean(line) not in boilerplate)
        for text in page_texts
    ]

    # Titre repéré dans le contenu : signet (table des matières) en priorité —
    # signal structurel bien plus fiable qu'une position de ligne — sinon
    # première ligne substantielle du corps (hors en-tête/pied de page).
    content_title_candidates: list[tuple[str, str]] = []
    try:
        for entry in reader.outline or []:
            bookmark_title = getattr(entry, "title", None)
            if bookmark_title:
                content_title_candidates.append((str(bookmark_title), "signet"))
                break
    except Exception:
        pass
    for text in body_texts:
        for line in text.splitlines():
            cleaned = _clean(line)
            if cleaned:
                content_title_candidates.append((cleaned, "position"))
                break
        if len(content_title_candidates) and content_title_candidates[-1][1] == "position":
            break
    content_title = _pick_content_title(content_title_candidates, path.name)
    if content_title:
        fields["Titre (contenu)"] = content_title

    keywords = _clean(meta.get("/Keywords"))
    if keywords:
        fields["Mots-clés"] = keywords

    theme = _recurring_keywords("\n".join(body_texts))
    if theme:
        fields["Mots du contenu"] = theme

    subject = _clean(meta.get("/Subject"))
    if subject:
        fields["Sujet"] = subject

    if boilerplate:
        fields["En-tête/pied de page"] = " · ".join(sorted(boilerplate, key=len)[:2])

    # Extrait verbatim (corps nettoyé, rejeté si encodage de police illisible).
    for text in body_texts:
        snippet = _plausible_snippet(text)
        if snippet:
            fields["Extrait"] = snippet
            break

    return fields


def _docx_header_footer_text(doc) -> str:
    """Texte réel d'en-tête/pied de page — API directe (`section.header`/
    `.footer`), pas d'inférence nécessaire contrairement au PDF."""
    section = doc.sections[0]
    parts = []
    for part in (section.header, section.footer):
        text = _clean(" ".join(p.text for p in part.paragraphs))
        if text:
            parts.append(text)
    return " · ".join(parts)


def _extract_docx(path: Path) -> dict[str, str]:
    from docx import Document

    doc = Document(str(path))
    fields: dict[str, str] = {}
    props = doc.core_properties

    title = _clean(props.title)
    if title and not _is_noise_title(title, path.name):
        fields["Titre interne"] = title

    paragraphs = doc.paragraphs
    total = len(paragraphs)
    # Titre structurel : les 20 premiers paragraphes (un titre/intitulé de
    # style "Heading"/"Title" apparaît toujours en tête de document).
    title_candidates = []
    for i in range(min(20, total)):
        style = paragraphs[i].style
        title_candidates.append((paragraphs[i].text, style.name if style else ""))
    content_title = _pick_content_title(title_candidates, path.name)
    if content_title:
        fields["Titre (contenu)"] = content_title

    keywords = _clean(props.keywords)
    if keywords:
        fields["Mots-clés"] = keywords

    # Thème/extrait : paragraphes début/milieu/fin (+ fenêtre) — le thème
    # peut se trouver n'importe où dans un document long.
    anchors = _sample_indices(total)
    sample_idx = _windowed_indices(anchors, total, _SAMPLE_WINDOW)
    sampled_text = "\n".join(paragraphs[i].text for i in sample_idx)

    theme = _recurring_keywords(sampled_text)
    if theme:
        fields["Mots du contenu"] = theme

    subject = _clean(props.subject)
    if subject:
        fields["Sujet"] = subject

    header_footer = _docx_header_footer_text(doc)
    if header_footer:
        fields["En-tête/pied de page"] = header_footer

    snippet = _plausible_snippet(sampled_text)
    if snippet:
        fields["Extrait"] = snippet
    return fields


_DEFAULT_SHEET_NAME_RE = re.compile(r"^(feuil|feuille|sheet)\d*$", re.IGNORECASE)


def _xlsx_sampled_text(ws, window: int = _SAMPLE_WINDOW) -> str:
    """Valeurs textuelles de cellules autour de lignes début/milieu/fin, sans
    charger toute la feuille (bornes via `ws.max_row`, toujours read_only)."""
    max_row = ws.max_row or 0
    texts: list[str] = []
    for anchor in _sample_indices(max_row, want=3):
        row_1based = anchor + 1
        lo = max(1, row_1based - window)
        hi = min(max_row, row_1based + window)
        for row in ws.iter_rows(min_row=lo, max_row=hi, values_only=True):
            for value in row:
                if isinstance(value, str):
                    texts.append(value)
    return "\n".join(texts)


def _extract_xlsx(path: Path) -> dict[str, str]:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    fields: dict[str, str] = {}
    props = wb.properties

    title = _clean(props.title)
    if title and not _is_noise_title(title, path.name):
        fields["Titre interne"] = title

    keywords = _clean(props.keywords)
    if keywords:
        fields["Mots-clés"] = keywords

    # Pas de « Titre (contenu) » pour le tabulaire : une feuille de calcul n'a
    # pas de notion de titre dans le contenu — limitation assumée. Pas
    # d'« En-tête/pied de page » non plus : `oddHeader`/`oddFooter` ne sont
    # pas exposés par `ReadOnlyWorksheet` (mode toujours utilisé ici pour ne
    # jamais charger une feuille entière en mémoire) — même arbitrage.
    ws = wb.active
    if ws is not None:
        theme = _recurring_keywords(_xlsx_sampled_text(ws))
        if theme:
            fields["Mots du contenu"] = theme

    subject = _clean(props.subject)
    if subject:
        fields["Sujet"] = subject

    # Noms de feuille par défaut écartés (Feuil1/Sheet1…) : n'apportent rien.
    sheets = [
        s for s in wb.sheetnames
        if _clean(s) and not _DEFAULT_SHEET_NAME_RE.match(s.strip())
    ]
    if sheets:
        fields["Feuilles"] = ", ".join(sheets[:8])
    wb.close()
    return fields


def _extract_pptx(path: Path) -> dict[str, str]:
    from pptx import Presentation
    from pptx.enum.shapes import PP_PLACEHOLDER

    prs = Presentation(str(path))
    fields: dict[str, str] = {}
    props = prs.core_properties

    title = _clean(props.title)
    if title and not _is_noise_title(title, path.name):
        fields["Titre interne"] = title

    slides = prs.slides
    sample_idx = _sample_indices(len(slides))

    # Titre repéré dans le contenu : placeholder de titre — signal structurel
    # direct (fiable), pas une position devinée.
    content_title = ""
    for i in sample_idx:
        title_shape = slides[i].shapes.title
        if title_shape is None:
            continue
        candidate = _clean(title_shape.text)
        if candidate and not _is_noise_title(candidate, path.name):
            content_title = candidate
            break
    if content_title:
        fields["Titre (contenu)"] = content_title

    keywords = _clean(props.keywords)
    if keywords:
        fields["Mots-clés"] = keywords

    # Thème/extrait : slides début/milieu/fin plutôt que la seule première
    # slide porteuse de contenu.
    sampled_texts = []
    footer_text = ""
    for i in sample_idx:
        slide = slides[i]
        shape_texts = [_clean(sh.text) for sh in slide.shapes if sh.has_text_frame]
        sampled_texts.append(" ".join(t for t in shape_texts if t))
        if not footer_text:
            for shape in slide.placeholders:
                try:
                    if shape.placeholder_format.type == PP_PLACEHOLDER.FOOTER:
                        candidate = _clean(shape.text)
                        if candidate:
                            footer_text = candidate
                            break
                except Exception:
                    continue

    theme = _recurring_keywords(" ".join(sampled_texts))
    if theme:
        fields["Mots du contenu"] = theme

    subject = _clean(props.subject)
    if subject:
        fields["Sujet"] = subject

    # En-tête/pied de page : best-effort (absent sur la plupart des decks).
    if footer_text:
        fields["En-tête/pied de page"] = footer_text

    for text in sampled_texts:
        snippet = _plausible_snippet(text)
        if snippet:
            fields["Extrait"] = snippet
            break
    return fields


_EXTRACTORS = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".xlsx": _extract_xlsx,
    ".pptx": _extract_pptx,
}


def _format_description(fields: dict[str, str], max_chars: int) -> str:
    """Assemble les champs en une ligne compacte `Label : valeur · …`."""
    parts = [f"{label} : {val}" for label, val in fields.items() if val]
    desc = " · ".join(parts)
    if len(desc) > max_chars:
        desc = desc[: max_chars - 1].rstrip() + "…"
    return desc


# ── API publique ─────────────────────────────────────────────────────────────

def enrich_descriptions(
    df: pd.DataFrame,
    source_root: Path,
    *,
    overwrite: bool = False,
    max_chars: int = DEFAULT_MAX_CHARS,
    on_progress=None,
) -> tuple[pd.DataFrame, EnrichReport]:
    """Remplit `Content.Description` des Items à partir des binaires.

    - Ne traite que les lignes `Content.DescriptionLevel == "Item"`.
    - Ne touche pas aux descriptions déjà renseignées, sauf `overwrite=True`.
    - Tolérante : un fichier illisible/absent n'interrompt pas le lot ; il est
      comptabilisé dans le rapport.

    Retourne `(df_enrichi, rapport)`. Le DataFrame est une copie ; l'original
    n'est jamais modifié.
    """
    result = df.copy()
    report = EnrichReport()

    if "Content.Description" not in result.columns:
        result["Content.Description"] = ""
    desc = result["Content.Description"].fillna("")

    is_item = result.get("Content.DescriptionLevel") == "Item"
    item_idx = result.index[is_item] if is_item is not False else result.index[[]]

    for i in item_idx:
        report.total_items += 1
        file_value = _clean(result.at[i, "File"]) if "File" in result.columns else ""
        existing = _clean(desc.at[i])

        if existing and not overwrite:
            report.already_filled += 1
            continue
        if not file_value or file_value == ".":
            continue

        ext = Path(file_value).suffix.lower()
        if ext not in SUPPORTED_EXTENSIONS:
            report.unsupported += 1
            continue

        target = _resolve(source_root, file_value)
        if not target.is_file():
            report.missing += 1
            continue

        if on_progress:
            on_progress(file_value)

        try:
            fields = _EXTRACTORS[ext](target)
        except Exception as e:  # lecture/parse échoué : on isole, on continue
            report.errors.append(f"{file_value} : {type(e).__name__} {e}")
            continue

        description = _format_description(fields, max_chars)
        if not description:
            report.no_text += 1
            continue

        result.at[i, "Content.Description"] = description
        report.enriched += 1

    return result, report


def _sha256_file(path: Path) -> str:
    """Empreinte SHA-256 d'un binaire, lue par blocs (mémoire bornée)."""
    h = hashlib.sha256()
    with path.open("rb") as fh:
        for chunk in iter(lambda: fh.read(_HASH_CHUNK), b""):
            h.update(chunk)
    return h.hexdigest()


def fingerprint_files(
    df: pd.DataFrame,
    source_root: Path,
    *,
    overwrite: bool = False,
    on_progress=None,
) -> tuple[pd.DataFrame, FingerprintReport]:
    """Calcule l'empreinte SHA-256 des binaires des Items dans `FINGERPRINT_COLUMN`.

    Étape **locale, facultative, opt-in** (modèle `enrich`) : ouvre chaque
    binaire référencé par `File` pour en calculer un hash déterministe. Contrairement
    à l'extraction de texte, **toutes** les extensions sont hachées (un doublon strict
    peut être un .jpg, un .zip…). Aucun appel LLM ; rien ne quitte la machine.

    - Ne traite que les lignes `Content.DescriptionLevel == "Item"`.
    - Ne recalcule pas une empreinte déjà présente, sauf `overwrite=True`.
    - Tolérante : un binaire absent/illisible n'interrompt pas le lot (comptabilisé).

    Retourne `(df_avec_empreinte, rapport)`. Le DataFrame est une copie ; l'original
    n'est jamais modifié. La colonne d'empreinte est ensuite lue par
    `core.audit_scan` pour détecter les doublons stricts dans le digest AUD-001.
    """
    result = df.copy()
    report = FingerprintReport()

    if FINGERPRINT_COLUMN not in result.columns:
        result[FINGERPRINT_COLUMN] = ""
    existing_col = result[FINGERPRINT_COLUMN].fillna("")

    is_item = result.get("Content.DescriptionLevel") == "Item"
    item_idx = result.index[is_item] if is_item is not False else result.index[[]]

    for i in item_idx:
        report.total_items += 1
        file_value = _clean(result.at[i, "File"]) if "File" in result.columns else ""
        existing = _clean(existing_col.at[i])

        if existing and not overwrite:
            report.already_hashed += 1
            continue
        if not file_value or file_value == ".":
            report.skipped += 1
            continue

        target = _resolve(source_root, file_value)
        if not target.is_file():
            report.missing += 1
            continue

        if on_progress:
            on_progress(file_value)

        try:
            digest = _sha256_file(target)
        except Exception as e:  # lecture échouée : on isole, on continue
            report.errors.append(f"{file_value} : {type(e).__name__} {e}")
            continue

        result.at[i, FINGERPRINT_COLUMN] = digest
        report.hashed += 1

    return result, report
